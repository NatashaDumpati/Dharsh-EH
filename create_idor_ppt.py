#!/usr/bin/env python3
"""
Generates IDOR Vulnerability Lab presentation (PowerPoint).
Run: python create_idor_ppt.py
Add screenshots to screenshots/ folder - see screenshots/README.md for what to capture.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

# Slide dimensions
SLIDE_W = 13.333
SLIDE_H = 7.5

def add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    left = Inches(0.5)
    top = Inches(2.5)
    width = Inches(12.333)
    tf = slide.shapes.add_textbox(left, top, width, Inches(1.5))
    p = tf.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(5, 46, 22)  # emerald-900
    p.alignment = PP_ALIGN.CENTER
    if subtitle:
        tf2 = slide.shapes.add_textbox(left, Inches(4.2), width, Inches(1))
        p2 = tf2.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(24)
        p2.font.color.rgb = RGBColor(16, 185, 129)  # emerald-500
        p2.alignment = PP_ALIGN.CENTER
    return slide

def add_content_slide(prs, title, bullets, sub_bullets=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Title
    tf = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    p = tf.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(5, 46, 22)
    # Content
    tf2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.333), Inches(5.5))
    tf2.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf2.text_frame.paragraphs[i] if i < len(tf2.text_frame.paragraphs) else tf2.text_frame.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(18)
        p.space_after = Pt(12)
        if sub_bullets and i < len(sub_bullets):
            for sub in sub_bullets[i]:
                sp = tf2.text_frame.add_paragraph()
                sp.text = f"  ◦ {sub}"
                sp.font.size = Pt(14)
                sp.font.color.rgb = RGBColor(55, 65, 81)
                sp.space_after = Pt(6)
    return slide

def add_slide_with_image(prs, title, body_text, image_path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tf = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    p = tf.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(5, 46, 22)
    if body_text:
        tf2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(5), Inches(1.5))
        tf2.word_wrap = True
        p2 = tf2.text_frame.paragraphs[0]
        p2.text = body_text
        p2.font.size = Pt(14)
    if image_path and os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(0.5), Inches(2.8), width=Inches(12.333), height=Inches(4.2))
    else:
        placeholder = slide.shapes.add_textbox(Inches(2), Inches(3.5), Inches(9.333), Inches(2))
        ph = placeholder.text_frame.paragraphs[0]
        ph.text = f"[Add screenshot: {os.path.basename(image_path)}]"
        ph.font.size = Pt(18)
        ph.font.color.rgb = RGBColor(156, 163, 175)
        ph.alignment = PP_ALIGN.CENTER
    return slide

def main():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    screenshots_dir = os.path.join(os.path.dirname(__file__), "screenshots")

    # 1. Title
    add_title_slide(prs, "IDOR Vulnerability Lab", "Insecure Direct Object Reference - MediVault Demo")

    # 2. What is IDOR - Introduction
    add_content_slide(prs, "What is IDOR?", [
        "Insecure Direct Object Reference (IDOR) is an access control vulnerability",
        "Occurs when an application exposes a reference to an internal object (file, record, key) in the URL or request",
        "Attackers can manipulate these references to access unauthorized data",
        "OWASP Top 10 - Broken Access Control (A01:2021)",
        "Common in APIs, web apps, and mobile backends"
    ])

    # 3. IDOR - How it works
    add_content_slide(prs, "How IDOR Works", [
        "Application uses predictable identifiers (IDs) to fetch resources",
        "Example: /api/patient?id=1001 fetches patient 1001's record",
        "No server-side check: Does the logged-in user own this resource?",
        "Attacker changes id=1001 to id=1002 → accesses another user's data",
        "Often combined with: sequential IDs, no authorization checks, exposed APIs"
    ])

    # 4. Impact
    add_content_slide(prs, "Impact of IDOR", [
        "Data Breach: Access to PII, medical records, financial data",
        "Privacy Violation: View/modify other users' private information",
        "Compliance Fines: GDPR, HIPAA violations → heavy penalties",
        "Reputation Damage: Loss of trust, legal liability",
        "Escalation: From read-only to full account takeover in some cases"
    ])

    # 5. Real-World Examples
    add_content_slide(prs, "Real-World IDOR Examples", [
        "Facebook (2013): Access to private photos via user ID manipulation",
        "Uber: Drivers could view other drivers' earnings by changing IDs",
        "Healthcare portals: Patient records exposed via predictable document IDs",
        "Banking apps: Account statements of other users by altering account numbers",
        "Social media: Private messages, DMs accessible by ID enumeration"
    ])

    # 6. Mitigations
    add_content_slide(prs, "Mitigations", [
        "Implement proper authorization: Check if user has permission for requested resource",
        "Use indirect references: Map session tokens to internal IDs (don't expose real IDs)",
        "Apply principle of least privilege: Users only get access to their own data",
        "Validate ownership server-side: Never trust client-supplied IDs alone",
        "Use UUIDs instead of sequential IDs: Harder to enumerate",
        "Audit and test: Regular penetration testing, code reviews"
    ])

    # 7. Lab Demo - Overview
    add_content_slide(prs, "Lab Demo: MediVault Portal", [
        "MediVault: Simulated healthcare portal with intentional IDOR vulnerability",
        "10 patient accounts (alice@example.com ... jack@example.com), password: password123",
        "Vulnerable endpoints: ?id= for patient records, ?doc_id= for documents",
        "Demo flow: Login as Alice → change URL id to access Bob's records",
        "Admin easter egg: ?id=1337 reveals system logs"
    ])

    # 8. Demo Screenshot - Login
    add_slide_with_image(prs, "Demo Step 1: Login to MediVault",
        "Login as alice@example.com / password123",
        os.path.join(screenshots_dir, "01_login.png"))

    # 9. Demo Screenshot - Own records
    add_slide_with_image(prs, "Demo Step 2: View Own Records (Alice - ID 1001)",
        "Navigate to Medical Records. URL shows ?id=1001 (Alice's ID)",
        os.path.join(screenshots_dir, "02_own_records.png"))

    # 10. Demo Screenshot - IDOR exploit
    add_slide_with_image(prs, "Demo Step 3: IDOR Exploit - Change ID in URL",
        "Manually change ?id=1001 to ?id=1002 in address bar → Bob's private data exposed",
        os.path.join(screenshots_dir, "03_idor_exploit.png"))

    # 11. Demo Screenshot - Unauthorized access
    add_slide_with_image(prs, "Demo Step 4: Unauthorized Access Detected",
        "Header shows 'UNAUTHORIZED ACCESS DETECTED' - but data is still displayed (vulnerability)",
        os.path.join(screenshots_dir, "04_unauthorized_banner.png"))

    # 12. Demo Screenshot - Document IDOR
    add_slide_with_image(prs, "Demo Step 5: Document IDOR",
        "Change ?doc_id=D-1001-01 to D-1004-01 to view David's clinical documents",
        os.path.join(screenshots_dir, "05_document_idor.png"))

    # 13. Summary
    add_content_slide(prs, "Summary", [
        "IDOR = Exposing internal object references without proper authorization",
        "High impact: Data breach, privacy violations, compliance fines",
        "Mitigate: Server-side authorization, indirect references, UUIDs",
        "MediVault lab demonstrates IDOR in a realistic healthcare context",
        "Always validate: Does this user have permission to access this resource?"
    ])

    # 14. Thank you
    add_title_slide(prs, "Thank You", "Questions?")

    out_path = os.path.join(os.path.dirname(__file__), "IDOR_Vulnerability_Lab_Demo.pptx")
    prs.save(out_path)
    print(f"Presentation saved: {out_path}")
    print("Add screenshots to screenshots/ folder and re-run to embed them.")
    print("See screenshots/README.md for capture instructions.")

if __name__ == "__main__":
    main()
